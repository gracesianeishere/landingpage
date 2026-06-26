from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Dimensions: A4 landscape-style, but we'll do portrait A4
# A4 portrait: 210mm x 297mm = 7.48" x 10.58"
W = Inches(7.48)
H = Inches(10.58)

RED       = RGBColor(0x8B, 0x00, 0x00)
INK       = RGBColor(0x0D, 0x0D, 0x0D)
INK2      = RGBColor(0x3A, 0x3A, 0x3A)
INK3      = RGBColor(0x6B, 0x6B, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CREAM     = RGBColor(0xF8, 0xF5, 0xF1)
RULE      = RGBColor(0xD9, 0xD0, 0xC8)
DARK      = RGBColor(0x0D, 0x0D, 0x0D)

MG = Inches(0.6)   # margin

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper functions ────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, color=INK,
             align=PP_ALIGN.LEFT, italic=False,
             font_name="Calibri", wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_text_multiline(slide, lines, x, y, w, h,
                        font_size=12, bold=False, color=INK,
                        align=PP_ALIGN.LEFT, font_name="Calibri",
                        spacing_after=6):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_label(slide, text, x, y, w=Inches(4)):
    return add_text(slide, text, x, y, w, Pt(16),
                    font_size=8, bold=True, color=RED,
                    font_name="Calibri")

def page_footer(slide, page_num):
    # footer rule
    add_rect(slide, MG, H - Inches(0.45), W - MG*2, Pt(1), fill=RULE)
    add_text(slide, "EDTAN FOR FUN", MG, H - Inches(0.42), Inches(2), Inches(0.25),
             font_size=8, bold=True, color=INK3, font_name="Calibri")
    add_text(slide, str(page_num).zfill(2), W - MG - Inches(0.3), H - Inches(0.42),
             Inches(0.3), Inches(0.25),
             font_size=8, color=INK3, align=PP_ALIGN.RIGHT, font_name="Calibri")

def bullet_item(slide, text, x, y, w, font_size=11, color=INK2):
    # dot
    dot_size = Pt(5)
    add_rect(slide, x, y + Pt(font_size * 0.45), dot_size, dot_size, fill=RED)
    # text
    add_text(slide, text, x + Inches(0.18), y, w - Inches(0.18), Inches(0.28),
             font_size=font_size, color=color, font_name="Calibri")

# ════════════════════════════════════════════════════════════
# PAGE 1 — COVER
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full red background
add_rect(slide, 0, 0, W, H, fill=RED)

# Top bar
add_rect(slide, 0, 0, W, Inches(0.65), fill=RGBColor(0x7A, 0x00, 0x00))
add_text(slide, "BRAND PROFILE", MG, Inches(0.18), Inches(3), Inches(0.3),
         font_size=8, bold=True, color=RGBColor(0xCC, 0x88, 0x88), font_name="Calibri")
add_text(slide, "2025", W - MG - Inches(0.5), Inches(0.18), Inches(0.5), Inches(0.3),
         font_size=8, color=RGBColor(0xCC, 0x88, 0x88), align=PP_ALIGN.RIGHT, font_name="Calibri")

# Logo area — symbol block + EDTAN text
logo_y = Inches(3.8)
# Symbol block background
add_rect(slide, MG, logo_y, Inches(0.55), Inches(0.85),
         fill=RGBColor(0xA0, 0x10, 0x10))
# Circle symbol (approximated as small rect)
add_rect(slide, MG + Inches(0.14), logo_y + Inches(0.08), Inches(0.27), Inches(0.18),
         fill=RGBColor(0xFF, 0xFF, 0xFF))
# Triangle row
add_rect(slide, MG + Inches(0.14), logo_y + Inches(0.33), Inches(0.27), Inches(0.16),
         fill=RGBColor(0xFF, 0xFF, 0xFF))
# Plus row
add_rect(slide, MG + Inches(0.14), logo_y + Inches(0.57), Inches(0.27), Inches(0.16),
         fill=RGBColor(0xFF, 0xFF, 0xFF))

# EDTAN
add_text(slide, "EDTAN", MG + Inches(0.65), logo_y - Inches(0.1), Inches(4), Inches(0.9),
         font_size=64, bold=True, color=WHITE, font_name="Calibri")
# FOR FUN
add_text(slide, "FOR FUN", MG + Inches(1.5), logo_y + Inches(0.72), Inches(3), Inches(0.4),
         font_size=20, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Calibri")

# Divider line
add_rect(slide, MG, Inches(5.1), Inches(0.6), Pt(2),
         fill=RGBColor(0xAA, 0x44, 0x44))

# Tagline
add_text(slide, "SOUTHEAST ASIA'S DESIGNER TOY & COLLECTIBLE PLATFORM",
         MG, Inches(5.25), W - MG*2, Inches(0.35),
         font_size=8, color=RGBColor(0xCC, 0xAA, 0xAA), font_name="Calibri")

# Headline
add_text_multiline(slide,
    ["Where Collectors", "Meet Creators."],
    MG, Inches(5.75), W - MG*2, Inches(1.8),
    font_size=46, bold=True, color=WHITE, font_name="Calibri", spacing_after=0)

# Sub
add_text(slide,
    "A media and community platform dedicated to designer toys, sofubi, art toys, and collectible culture across Southeast Asia and beyond.",
    MG, Inches(7.55), Inches(4.5), Inches(0.9),
    font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Calibri")

# Bottom bar
add_rect(slide, 0, H - Inches(0.55), W, Pt(1), fill=RGBColor(0xAA, 0x44, 0x44))
add_text(slide, "CONFIDENTIAL · BRAND PROFILE 2025",
         MG, H - Inches(0.5), Inches(3), Inches(0.35),
         font_size=8, color=RGBColor(0x99, 0x55, 0x55), font_name="Calibri")
add_text(slide, "EdTan For Fun",
         W - MG - Inches(1.5), H - Inches(0.5), Inches(1.5), Inches(0.35),
         font_size=8, color=RGBColor(0x99, 0x55, 0x55), align=PP_ALIGN.RIGHT, font_name="Calibri")

# ════════════════════════════════════════════════════════════
# PAGE 2 — ABOUT + VISION + MISSION
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "ABOUT", MG, Inches(0.55))

# Pull quote (left col)
col1_x = MG
col2_x = MG + Inches(3.2)
col_w  = Inches(3.0)

add_text_multiline(slide,
    ["A platform built for", "stories as much as", "collectibles."],
    col1_x, Inches(0.95), col_w, Inches(1.6),
    font_size=22, bold=True, color=INK, font_name="Calibri", spacing_after=2)

# Accent on "stories" and "collectibles" — we'll add a note in red below
add_text(slide, "(Key words in red in final design)", col1_x, Inches(2.45),
         col_w, Inches(0.25), font_size=7, color=INK3, font_name="Calibri")

# Body copy (right col)
body = (
    "EdTan For Fun is a media and community platform dedicated to designer toys, "
    "sofubi, art toys, and collectible culture.\n\n"
    "Through editorial storytelling, artist spotlights, event coverage, and community-driven "
    "initiatives, EdTan connects collectors, creators, and brands within a growing creative ecosystem.\n\n"
    "More than showcasing collectibles, EdTan highlights the stories behind the artists, craftsmanship, "
    "and creative process that bring each piece to life — while helping emerging creators gain greater "
    "visibility within the global designer toy scene."
)
txBox = slide.shapes.add_textbox(col2_x, Inches(0.95), col_w, Inches(1.9))
txBox.word_wrap = True
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = body
run.font.size = Pt(10)
run.font.color.rgb = INK2
run.font.name = "Calibri"

# Rule
add_rect(slide, MG, Inches(2.95), W - MG*2, Pt(1), fill=RULE)

# Vision
add_label(slide, "OUR VISION", MG, Inches(3.1))
add_rect(slide, MG, Inches(3.35), Pt(3), Inches(0.75), fill=RED)
add_text(slide,
    "To become one of Southeast Asia's leading platforms for designer toys and collectible culture — "
    "connecting collectors, creators, communities, and brands through meaningful creative initiatives.",
    MG + Inches(0.22), Inches(3.35), W - MG*2 - Inches(0.22), Inches(0.85),
    font_size=13, bold=True, color=INK, font_name="Calibri")
add_rect(slide, MG, Inches(3.3), W - MG*2, Inches(0.98), fill=CREAM)
add_rect(slide, MG, Inches(3.3), Pt(3), Inches(0.98), fill=RED)
# re-add text on top of cream box
add_text(slide,
    "To become one of Southeast Asia's leading platforms for designer toys and collectible culture — "
    "connecting collectors, creators, communities, and brands through meaningful creative initiatives.",
    MG + Inches(0.25), Inches(3.38), W - MG*2 - Inches(0.3), Inches(0.85),
    font_size=12, bold=False, color=INK, font_name="Calibri")

# Rule
add_rect(slide, MG, Inches(4.45), W - MG*2, Pt(1), fill=RULE)

# Mission
add_label(slide, "OUR MISSION", MG, Inches(4.6))
mission_items = ["Connecting collectors", "Empowering creators", "Growing collectible culture"]
my = Inches(4.88)
for item in mission_items:
    add_rect(slide, MG, my + Pt(8), Pt(7), Pt(7), fill=RED)
    add_text(slide, item, MG + Inches(0.22), my, W - MG*2, Inches(0.4),
             font_size=20, bold=True, color=INK, font_name="Calibri")
    my += Inches(0.47)

page_footer(slide, 2)

# ════════════════════════════════════════════════════════════
# PAGE 3 — WHAT WE DO
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "WHAT WE DO", MG, Inches(0.55))
add_text_multiline(slide, ["Three Pillars.", "One Ecosystem."],
    MG, Inches(0.82), W - MG*2, Inches(1.0),
    font_size=30, bold=True, color=INK, font_name="Calibri", spacing_after=0)
add_text(slide,
    "EdTan operates across three interconnected disciplines — editorial media, community building, and creative initiatives — each reinforcing the other.",
    MG, Inches(1.78), W - MG*2, Inches(0.5),
    font_size=10, color=INK3, font_name="Calibri")

# Rule
add_rect(slide, MG, Inches(2.38), W - MG*2, Pt(1), fill=RULE)

# 3 columns
col_w3 = (W - MG*2) / 3
cols = [
    {
        "num": "01 — EDITORIAL MEDIA",
        "title": "Original Content & Storytelling",
        "body": "We create editorial content that goes beyond product promotion — focusing on the artists, the craft, and the culture that define the designer toy world.",
        "items": ["Designer Toy Features", "Artist Spotlights", "Collection Highlights", "Industry News", "Event Coverage", "Product Reviews"]
    },
    {
        "num": "02 — COMMUNITY",
        "title": "Connecting Collectors & Creators",
        "body": "We foster an active community where collectors and creators connect, exchange ideas, and grow together through shared passion for the art form.",
        "items": ["EdTanVerse", "Collector Engagement", "Creator Discovery", "Community Discussions"]
    },
    {
        "num": "03 — CREATIVE INITIATIVES",
        "title": "Empowering Emerging Artists",
        "body": "EdTan continuously develops programs that encourage creativity and provide emerging artists and designers with platforms to reach wider audiences.",
        "items": ["Design Competitions", "Artist Collaborations", "Creator Showcases", "Workshops", "Exhibitions"]
    }
]

for i, col in enumerate(cols):
    cx = MG + col_w3 * i
    cy = Inches(2.5)
    # vertical rule between cols
    if i > 0:
        add_rect(slide, cx, Inches(2.38), Pt(1), Inches(5.5), fill=RULE)
    add_text(slide, col["num"], cx + Inches(0.1), cy, col_w3 - Inches(0.1), Inches(0.25),
             font_size=8, bold=True, color=RED, font_name="Calibri")
    add_text(slide, col["title"], cx + Inches(0.1), cy + Inches(0.3), col_w3 - Inches(0.15), Inches(0.55),
             font_size=13, bold=True, color=INK, font_name="Calibri")
    add_text(slide, col["body"], cx + Inches(0.1), cy + Inches(0.9), col_w3 - Inches(0.15), Inches(1.1),
             font_size=9.5, color=INK2, font_name="Calibri")
    iy = cy + Inches(2.1)
    for item in col["items"]:
        add_rect(slide, cx + Inches(0.1), iy + Pt(5), Pt(5), Pt(5), fill=RED)
        add_text(slide, item, cx + Inches(0.28), iy, col_w3 - Inches(0.3), Inches(0.28),
                 font_size=9.5, color=INK3, font_name="Calibri")
        iy += Inches(0.28)

# Content tags
add_rect(slide, MG, Inches(7.88), W - MG*2, Pt(1), fill=RULE)
add_label(slide, "CONTENT FOCUS AREAS", MG, Inches(8.02))
tags = ["Designer Toys", "Sofubi", "Art Toys", "Artist Spotlights", "Event Coverage", "Pop Culture", "Product Reviews"]
tx = MG
for tag in tags:
    tw = Inches(len(tag) * 0.085 + 0.22)
    add_rect(slide, tx, Inches(8.28), tw, Inches(0.28), line=RULE, line_w=Pt(1))
    add_text(slide, tag, tx + Inches(0.08), Inches(8.3), tw, Inches(0.24),
             font_size=8.5, bold=True, color=INK2, font_name="Calibri")
    tx += tw + Inches(0.1)
    if tx > W - MG - Inches(1.2):
        tx = MG

page_footer(slide, 3)

# ════════════════════════════════════════════════════════════
# PAGE 4 — AUDIENCE SNAPSHOT
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "AUDIENCE SNAPSHOT", MG, Inches(0.55))
add_text_multiline(slide, ["Engaged. Niche.", "Growing."],
    MG, Inches(0.82), W - MG*2, Inches(0.95),
    font_size=30, bold=True, color=INK, font_name="Calibri", spacing_after=0)
add_text(slide,
    "A highly targeted community with above-average engagement — every metric reflects a genuine, passionate audience rather than passive followers.",
    MG, Inches(1.72), W - MG*2, Inches(0.45),
    font_size=10, color=INK3, font_name="Calibri")

# Metrics grid
metrics = [
    ("1,214", "Instagram\nFollowers"),
    ("6.73%", "Average\nEngagement Rate"),
    ("4,821", "Average Views\nper Post"),
    ("15K+",  "Top Post\nReach"),
]
add_rect(slide, MG, Inches(2.3), W - MG*2, Inches(1.25), line=RULE, line_w=Pt(1))
mw = (W - MG*2) / 4
for i, (val, lbl) in enumerate(metrics):
    mx = MG + mw * i
    if i > 0:
        add_rect(slide, mx, Inches(2.3), Pt(1), Inches(1.25), fill=RULE)
    add_text(slide, val, mx + Inches(0.15), Inches(2.38), mw - Inches(0.15), Inches(0.65),
             font_size=28, bold=True, color=RED, font_name="Calibri")
    add_text(slide, lbl, mx + Inches(0.15), Inches(3.0), mw - Inches(0.15), Inches(0.45),
             font_size=8.5, color=INK3, font_name="Calibri")

# Engagement note
add_rect(slide, MG, Inches(3.72), W - MG*2, Inches(0.6), fill=CREAM)
add_text(slide,
    "A 6.73% engagement rate sits well above the typical Instagram benchmark of 1–3%, reflecting a highly active and invested community rather than passive follower counts.",
    MG + Inches(0.15), Inches(3.78), W - MG*2 - Inches(0.2), Inches(0.5),
    font_size=9.5, color=INK2, font_name="Calibri")

# Per-post bars
add_label(slide, "AVERAGE PERFORMANCE PER POST", MG, Inches(4.52))

bar_data = [
    ("Views",       4821, 4821),
    ("Engagements", 82,   4821),
    ("Likes",       77,   4821),
    ("Comments",    4,    4821),
]
bar_max_w = Inches(4.5)
bar_x     = MG + Inches(1.2)
by = Inches(4.82)
for label, val, maxv in bar_data:
    add_text(slide, label, MG, by, Inches(1.1), Inches(0.28),
             font_size=9.5, color=INK2, align=PP_ALIGN.RIGHT, font_name="Calibri")
    # track
    add_rect(slide, bar_x, by + Pt(6), bar_max_w, Pt(8), fill=RGBColor(0xED, 0xE8, 0xE2))
    # fill — use a minimum visible width
    fill_w = max(bar_max_w * val / maxv, Inches(0.04))
    add_rect(slide, bar_x, by + Pt(6), fill_w, Pt(8), fill=RED)
    # value
    add_text(slide, f"{val:,}", bar_x + bar_max_w + Inches(0.1), by, Inches(0.6), Inches(0.28),
             font_size=10, bold=True, color=INK, font_name="Calibri")
    by += Inches(0.42)

page_footer(slide, 4)

# ════════════════════════════════════════════════════════════
# PAGE 5 — DEMOGRAPHICS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "AUDIENCE DEMOGRAPHICS", MG, Inches(0.55))
add_text_multiline(slide, ["Who Is", "Our Audience?"],
    MG, Inches(0.82), W - MG*2, Inches(0.95),
    font_size=30, bold=True, color=INK, font_name="Calibri", spacing_after=0)

# Gender col
gc_x = MG
gc_w = Inches(3.0)
add_text(slide, "GENDER SPLIT", gc_x, Inches(1.95), gc_w, Inches(0.25),
         font_size=8, bold=True, color=INK3, font_name="Calibri")

# Donut approximation: two arcs as stacked rects with labels
add_rect(slide, gc_x, Inches(2.28), Inches(1.1), Inches(1.1),
         fill=RGBColor(0xED, 0xE8, 0xE2))
add_rect(slide, gc_x + Inches(0.05), Inches(2.33), Inches(1.0), Inches(0.6), fill=RED)
add_text(slide, "85%\nMale", gc_x + Inches(1.25), Inches(2.38), Inches(0.8), Inches(0.5),
         font_size=18, bold=True, color=RED, font_name="Calibri")

legend_items = [("Male", "85.2%", RED), ("Female", "14.8%", RGBColor(0xD9, 0xC8, 0xC8))]
ly = Inches(3.52)
for lbl, pct, clr in legend_items:
    add_rect(slide, gc_x, ly + Pt(3), Pt(10), Pt(10), fill=clr)
    add_text(slide, f"{lbl}  {pct}", gc_x + Inches(0.2), ly, Inches(1.5), Inches(0.28),
             font_size=10, color=INK2, font_name="Calibri")
    ly += Inches(0.32)

add_rect(slide, gc_x, Inches(4.0), gc_w, Inches(0.55), fill=CREAM)
add_text(slide, "Predominantly male audience consistent with global designer toy collector communities.",
         gc_x + Inches(0.1), Inches(4.04), gc_w - Inches(0.15), Inches(0.48),
         font_size=8.5, color=INK2, font_name="Calibri")

# Age col
ac_x = MG + Inches(3.3)
ac_w = Inches(3.55)
add_text(slide, "AGE DISTRIBUTION", ac_x, Inches(1.95), ac_w, Inches(0.25),
         font_size=8, bold=True, color=INK3, font_name="Calibri")

age_data = [
    ("35–44", 39.1, RED),
    ("25–34", 34.8, RGBColor(0xC4, 0x70, 0x70)),
    ("45–54", 16.3, RGBColor(0xD9, 0xA0, 0xA0)),
    ("18–24", 6.0,  RGBColor(0xE8, 0xC4, 0xC4)),
    ("Others", 3.8, RGBColor(0xF0, 0xD8, 0xD8)),
]
ay = Inches(2.28)
bar_w_max = Inches(2.2)
for age, pct, clr in age_data:
    add_text(slide, age, ac_x, ay, Inches(0.55), Inches(0.28),
             font_size=9.5, color=INK2, font_name="Calibri")
    add_rect(slide, ac_x + Inches(0.6), ay + Pt(7), bar_w_max, Pt(7),
             fill=RGBColor(0xED, 0xE8, 0xE2))
    add_rect(slide, ac_x + Inches(0.6), ay + Pt(7),
             max(bar_w_max * pct / 39.1, Inches(0.05)), Pt(7), fill=clr)
    add_text(slide, f"{pct}%", ac_x + Inches(0.6) + bar_w_max + Inches(0.08), ay,
             Inches(0.45), Inches(0.28), font_size=9.5, bold=True, color=INK, font_name="Calibri")
    ay += Inches(0.35)

add_rect(slide, ac_x, Inches(4.0), ac_w, Inches(0.55), fill=CREAM)
add_text(slide, "More than 90% of our audience falls within the 25–54 age group — a mature, financially established community with deep investment in collectible culture.",
         ac_x + Inches(0.1), Inches(4.04), ac_w - Inches(0.15), Inches(0.48),
         font_size=8.5, color=INK2, font_name="Calibri")

# Geographic
add_rect(slide, MG, Inches(4.75), W - MG*2, Pt(1), fill=RULE)
add_text(slide, "GEOGRAPHIC REACH — TOP MARKETS", MG, Inches(4.9), W - MG*2, Inches(0.25),
         font_size=8, bold=True, color=INK3, font_name="Calibri")

markets = [("🇮🇩", "Indonesia"), ("🇵🇭", "Philippines"), ("🇹🇼", "Taiwan"), ("🇭🇰", "Hong Kong"), ("🇹🇭", "Thailand")]
gw = (W - MG*2) / 5
for i, (flag, name) in enumerate(markets):
    gx = MG + gw * i
    add_rect(slide, gx + Inches(0.05), Inches(5.1), gw - Inches(0.1), Inches(0.75), line=RULE, line_w=Pt(1))
    add_text(slide, flag, gx + Inches(0.05), Inches(5.12), gw - Inches(0.1), Inches(0.35),
             font_size=18, align=PP_ALIGN.CENTER, font_name="Segoe UI Emoji")
    add_text(slide, name, gx + Inches(0.05), Inches(5.48), gw - Inches(0.1), Inches(0.28),
             font_size=8.5, bold=True, color=INK2, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(slide,
    "As EdTan publishes primarily in English, our content reaches collectors and creators beyond Indonesia — connecting us with a broader international audience across the Asia-Pacific region.",
    MG, Inches(6.02), W - MG*2, Inches(0.45),
    font_size=9.5, color=INK3, font_name="Calibri")

page_footer(slide, 5)

# ════════════════════════════════════════════════════════════
# PAGE 6 — WHY PARTNER
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

# Dark header block
header_h = Inches(3.2)
add_rect(slide, 0, 0, W, header_h, fill=DARK)
add_text(slide, "PARTNERSHIP", MG, Inches(0.55), Inches(3), Inches(0.25),
         font_size=8, bold=True, color=RGBColor(0xB8, 0x44, 0x44), font_name="Calibri")
add_text_multiline(slide, ["Why Partner", "With EdTan?"],
    MG, Inches(0.9), W - MG*2, Inches(1.3),
    font_size=34, bold=True, color=WHITE, font_name="Calibri", spacing_after=0)
add_text(slide,
    "Partnering with EdTan means reaching a highly targeted audience that genuinely appreciates creativity, craftsmanship, and collectible culture — through authentic editorial storytelling rather than conventional advertising.",
    MG, Inches(2.15), Inches(5.5), Inches(0.85),
    font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), font_name="Calibri")

# Reasons
reasons = [
    ("01", "Highly Targeted Collectible Audience",
     "A community built around genuine passion — not casual interest. Every follower has a reason to be here."),
    ("02", "Strong Niche Engagement",
     "A 6.73% average engagement rate demonstrates a community that actively participates, not just passively scrolls."),
    ("03", "Editorial-First Storytelling",
     "Content that feels authentic because it is — collaborations are shaped around stories, not ad slots."),
    ("04", "English-First, Regionally Rooted",
     "English-language content that naturally connects brands to audiences across Southeast Asia and beyond."),
    ("05", "Trusted Platform for the Community",
     "A recognized voice within the designer toy and collectible enthusiast space — credibility built over time."),
]

ry = Inches(3.3)
for num, title, body in reasons:
    add_rect(slide, MG, ry + Inches(0.32), W - MG*2, Pt(1), fill=RULE)
    add_text(slide, num, MG, ry + Inches(0.05), Inches(0.35), Inches(0.25),
             font_size=10, bold=True, color=RED, font_name="Calibri")
    add_text(slide, title, MG + Inches(0.42), ry + Inches(0.02), W - MG*2 - Inches(0.42), Inches(0.25),
             font_size=12, bold=True, color=INK, font_name="Calibri")
    add_text(slide, body, MG + Inches(0.42), ry + Inches(0.24), W - MG*2 - Inches(0.42), Inches(0.3),
             font_size=9.5, color=INK2, font_name="Calibri")
    ry += Inches(0.68)

page_footer(slide, 6)

# ════════════════════════════════════════════════════════════
# PAGE 7 — COLLABORATIONS & HIGHLIGHTS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "PREVIOUS COLLABORATIONS", MG, Inches(0.55))
add_text_multiline(slide, ["Brands We've", "Worked With"],
    MG, Inches(0.82), W - MG*2, Inches(0.95),
    font_size=30, bold=True, color=INK, font_name="Calibri", spacing_after=0)
add_text(slide,
    "EdTan has collaborated with and featured various brands and key industry players across the designer toy ecosystem.",
    MG, Inches(1.72), W - MG*2, Inches(0.38),
    font_size=10, color=INK3, font_name="Calibri")

# Collab grid
add_rect(slide, MG, Inches(2.22), W - MG*2, Pt(1), fill=RULE)
brands = ["Amazing Toy Show", "Unbox Industries", "ThePrjctPlayground", "ERON Collects"]
bw = (W - MG*2) / 2
for i, brand in enumerate(brands):
    bx = MG + (bw * (i % 2))
    by_b = Inches(2.28) + Inches(0.72) * (i // 2)
    add_rect(slide, bx, by_b, bw, Inches(0.72), line=RULE, line_w=Pt(1))
    add_rect(slide, bx + Inches(0.25), by_b + Inches(0.32), Pt(8), Pt(8), fill=RED)
    add_text(slide, brand, bx + Inches(0.48), by_b + Inches(0.24), bw - Inches(0.55), Inches(0.3),
             font_size=14, bold=True, color=INK, font_name="Calibri")

# Highlights
add_rect(slide, MG, Inches(3.85), W - MG*2, Pt(1), fill=RULE)
add_text(slide, "BRAND HIGHLIGHTS", MG, Inches(4.0), W - MG*2, Inches(0.25),
         font_size=8, bold=True, color=INK3, font_name="Calibri")

highlights = [
    "Built EdTanVerse — a growing community space dedicated to designer toy and collectible enthusiasts.",
    "Event coverage consistently generates strong engagement, with top-performing posts reaching 10K–15K+ views.",
    "Recognized for delivering editorial content focused on artists, collectible stories, and industry events.",
    "Continues to expand its role beyond media through community initiatives and creator-focused programs.",
]
hw = (W - MG*2 - Inches(0.15)) / 2
hcard_h = Inches(1.0)
for i, h in enumerate(highlights):
    hx = MG + (hw + Inches(0.15)) * (i % 2)
    hy = Inches(4.32) + (hcard_h + Inches(0.12)) * (i // 2)
    add_rect(slide, hx, hy, hw, hcard_h, fill=CREAM)
    add_text(slide, h, hx + Inches(0.15), hy + Inches(0.12), hw - Inches(0.25), hcard_h - Inches(0.18),
             font_size=9.5, color=INK2, font_name="Calibri")

page_footer(slide, 7)

# ════════════════════════════════════════════════════════════
# PAGE 8 — IDEAL PARTNERS + CLOSING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=WHITE)

add_label(slide, "LET'S WORK TOGETHER", MG, Inches(0.55))
add_text_multiline(slide, ["Ideal", "Collaboration Partners"],
    MG, Inches(0.82), W - MG*2, Inches(0.95),
    font_size=30, bold=True, color=INK, font_name="Calibri", spacing_after=0)
add_text(slide,
    "EdTan is open to partnerships with any brand, studio, creator, or organization that shares a genuine connection to the designer toy and creative collectible ecosystem.",
    MG, Inches(1.72), W - MG*2, Inches(0.45),
    font_size=10, color=INK3, font_name="Calibri")

# Partner types grid
add_rect(slide, MG, Inches(2.32), W - MG*2, Pt(1), fill=RULE)
partner_types = [
    "Designer Toy Brands", "Sofubi Studios",
    "Art Toy Creators", "Toy Conventions & Events",
    "Collectible Retailers", "Creative Lifestyle Brands",
    "Art Galleries", "Pop Culture & Creative Communities",
]
pw = (W - MG*2) / 2
ph = Inches(0.52)
for i, pt_name in enumerate(partner_types):
    px = MG + pw * (i % 2)
    py = Inches(2.32) + ph * (i // 2)
    add_rect(slide, px, py, pw, ph, line=RULE, line_w=Pt(1))
    add_rect(slide, px + Inches(0.2), py + Inches(0.22), Pt(7), Pt(7), fill=RED)
    add_text(slide, pt_name, px + Inches(0.42), py + Inches(0.14), pw - Inches(0.5), Inches(0.28),
             font_size=11, bold=True, color=INK, font_name="Calibri")

# Closing panel
cl_y = Inches(6.55)
cl_h = Inches(3.42)
add_rect(slide, MG, cl_y, W - MG*2, cl_h, fill=RED)

add_text(slide, "LOOKING FORWARD",
         MG + Inches(0.4), cl_y + Inches(0.3), W - MG*2, Inches(0.25),
         font_size=8, bold=True, color=RGBColor(0xCC, 0x88, 0x88), font_name="Calibri")

mission_lines = ["Connecting collectors", "Empowering creators", "Growing collectible culture"]
my2 = cl_y + Inches(0.65)
for line in mission_lines:
    add_rect(slide, MG + Inches(0.4), my2 + Pt(10), Pt(7), Pt(7), fill=RGBColor(0xCC, 0x88, 0x88))
    add_text(slide, line, MG + Inches(0.65), my2, W - MG*2 - Inches(0.6), Inches(0.42),
             font_size=22, bold=True, color=WHITE, font_name="Calibri")
    my2 += Inches(0.45)

add_text(slide,
    "EdTan For Fun is committed to growing beyond a media platform — developing initiatives that empower creators, strengthen communities, and expand the designer toy ecosystem across Southeast Asia and beyond.",
    MG + Inches(0.4), cl_y + Inches(2.1), W - MG*2 - Inches(0.5), Inches(0.85),
    font_size=10, color=RGBColor(0xCC, 0xAA, 0xAA), font_name="Calibri")

page_footer(slide, 8)

# ── Save ──────────────────────────────────────────────────
out = "/home/user/landingpage/EdTan_For_Fun_Brand_Profile.pptx"
prs.save(out)
print(f"Saved: {out}")
